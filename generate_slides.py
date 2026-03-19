"""
Generate Capstone Presentation Slides (.pptx)
EDUC 9510 – Masters Capstone Action Research
Fitchburg State University, Spring 2026

Generates: capstone_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette (matches app theme) ──────────────────────────────
BLUE       = RGBColor(0x1E, 0x88, 0xE5)   # #1E88E5  primary accent
DARK_BLUE  = RGBColor(0x0D, 0x47, 0xA1)   # #0D47A1  darker heading
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x1A)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
GREEN      = RGBColor(0x09, 0xAB, 0x3B)   # #09AB3B  app green
YELLOW_C   = RGBColor(0xFA, 0xCA, 0x2B)   # app yellow
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)   # app orange
RED        = RGBColor(0xFF, 0x2B, 0x2B)   # app red

prs = Presentation()
prs.slide_width  = Inches(13.333)   # widescreen 16:9
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────────────

def add_bg(slide, color=WHITE):
    """Set solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    """Add a text box with styled text."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txbox


def add_bullet_frame(slide, left, top, width, height, bullets, font_size=18,
                     color=DARK_TEXT, bold_first=False, spacing=10,
                     font_name="Calibri"):
    """Add a text frame with bullet points."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(spacing)
        p.level = 0
        if bold_first and i == 0:
            p.font.bold = True
    return txbox


def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def section_header_bar(slide, text):
    """Add a coloured bar at the top with a section title."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
                text, font_size=32, bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT)


def add_color_swatch(slide, left, top, color, label):
    """Draw a small colour swatch + label (for the scoring legend)."""
    add_rect(slide, left, top, Inches(0.35), Inches(0.35), color)
    add_textbox(slide, left + Inches(0.45), top - Inches(0.02),
                Inches(1.6), Inches(0.4), label,
                font_size=14, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)

# Blue accent bar at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.25), BLUE)

# Title
add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(2),
            "Impact of UFLI Assessment and Reteaching Sequence\n"
            "on Kindergarten Letter Identification and Letter Sounds",
            font_size=36, bold=True, color=DARK_BLUE,
            alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Subtitle / author info
add_textbox(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(2.5),
            "Masters Capstone Action Research\n"
            "EDUC 9510  •  Fitchburg State University  •  Spring 2026",
            font_size=22, color=MED_GRAY, alignment=PP_ALIGN.CENTER,
            line_spacing=1.4)

# Bottom accent bar
add_rect(slide, Inches(0), SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), BLUE)

add_speaker_notes(slide,
    "Welcome and introduction. State your name, program, and briefly introduce "
    "the topic: using UFLI assessment and reteaching to support kindergarteners "
    "struggling with letter identification after Tier 1 instruction.")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 – Research Question
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Research Question")

# Quote-style box
add_rect(slide, Inches(1.2), Inches(1.8), Inches(10.9), Inches(3.5),
         LIGHT_GRAY, line_color=BLUE)

# Left accent line
add_rect(slide, Inches(1.2), Inches(1.8), Inches(0.12), Inches(3.5), BLUE)

add_textbox(slide, Inches(1.8), Inches(2.2), Inches(9.8), Inches(2.8),
            '"What is the impact of using the UFLI assessment and '
            'reteaching sequence with my Kindergarten students struggling '
            'to recognize letter identification and letter sounds after '
            'Tier 1 instruction?"',
            font_size=26, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
            line_spacing=1.5, font_name="Calibri")

# Context
add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1),
            "Target Population: Kindergarten students who have not met "
            "letter identification benchmarks following Tier 1 instruction",
            font_size=18, color=MED_GRAY, alignment=PP_ALIGN.LEFT)

add_speaker_notes(slide,
    "Read the research question aloud. Emphasize the focus on students who are "
    "struggling AFTER receiving standard Tier 1 instruction — this is about "
    "targeted intervention, not replacing core instruction.")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 – Operational Definitions
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Operational Definitions")

definitions = [
    ("Letter Recognition", "To identify a specific letter from the alphabet"),
    ("Letter Sounds", "To identify a specific letter and the sound it makes"),
    ("Letter-Sound Correspondence",
     "The relationship between a sound (phoneme) and the letter(s) (grapheme) that spell it"),
    ("Alphabetic Principle",
     "The understanding that spoken words are represented by letters in written language"),
    ("Literacy Skills", "Letter name identification and letter sound recognition"),
    ("Sensorimotor Experiences", "Hands-on ways to understand and learn"),
]

y = Inches(1.5)
for term, defn in definitions:
    # Term in bold blue
    add_textbox(slide, Inches(1), y, Inches(4), Inches(0.45),
                term, font_size=18, bold=True, color=BLUE)
    # Definition
    add_textbox(slide, Inches(4.8), y, Inches(7.5), Inches(0.45),
                defn, font_size=17, color=DARK_TEXT)
    y += Inches(0.65)

# Thin divider line across
add_rect(slide, Inches(1), Inches(1.35), Inches(11.3), Inches(0.02), BLUE)

add_speaker_notes(slide,
    "Briefly review key terms. Don't read every definition — highlight the most "
    "important ones (letter recognition, letter sounds, letter-sound correspondence) "
    "and note that these definitions guided your data collection and analysis.")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 – The Intervention
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "The Intervention")

# Left column
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
            "UFLI Assessment & Reteaching Sequence",
            font_size=24, bold=True, color=DARK_BLUE)

add_bullet_frame(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4),
    [
        "Systematic assessment of all 26 letters across three pillars",
        "Identifies specific letters each student has not yet mastered",
        "Targeted reteaching based on individual assessment data",
        "Progress monitored through repeated assessments over time",
        "Intervention delivered to students below benchmark after Tier 1",
    ],
    font_size=17, spacing=12)

# Right column – "Three Pillars" visual
# Pillar boxes
pillar_data = [
    ("ABC", "Uppercase\nRecognition", "Can the student identify\neach uppercase letter?"),
    ("abc", "Lowercase\nRecognition", "Can the student identify\neach lowercase letter?"),
    ("🔊", "Letter\nSounds", "Can the student produce\neach letter sound?"),
]

x_start = Inches(7.2)
for i, (icon, title, desc) in enumerate(pillar_data):
    x = x_start + Inches(i * 2.05)
    # Card background
    add_rect(slide, x, Inches(1.8), Inches(1.9), Inches(4.0), LIGHT_GRAY, line_color=BLUE)
    # Icon
    add_textbox(slide, x, Inches(2.0), Inches(1.9), Inches(0.7),
                icon, font_size=30, color=BLUE, alignment=PP_ALIGN.CENTER,
                bold=True)
    # Title
    add_textbox(slide, x + Inches(0.1), Inches(2.7), Inches(1.7), Inches(0.9),
                title, font_size=16, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER, line_spacing=1.2)
    # Description
    add_textbox(slide, x + Inches(0.1), Inches(3.7), Inches(1.7), Inches(1.0),
                desc, font_size=13, color=MED_GRAY,
                alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# "26 letters each → 78 total items" label
add_textbox(slide, Inches(7.2), Inches(6.0), Inches(6), Inches(0.5),
            "26 letters × 3 pillars = 78 items per assessment",
            font_size=15, color=MED_GRAY, alignment=PP_ALIGN.CENTER,
            bold=True)

add_speaker_notes(slide,
    "Explain the intervention: each student is assessed on all 26 letters across "
    "three categories — uppercase recognition, lowercase recognition, and letter sounds. "
    "This totals 78 data points per student per assessment. Reteaching targets the "
    "specific letters a student has NOT mastered.")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 – The Technology Tool: Letter ID Tracker
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "The Tool: Letter ID Tracker")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6),
            "A custom-built Streamlit web application for real-time letter identification tracking",
            font_size=20, color=MED_GRAY)

# Feature cards (2 rows × 3 cols)
features = [
    ("Live Assessment", "Interactive flashcard-based assessment\nwith sequential three-phase flow\n(Uppercase → Lowercase → Sounds)"),
    ("Student Scorecard", "At-a-glance metrics for every student\nwith color-coded performance\nindicators (green/yellow/orange/red)"),
    ("Individual Progress", "Track each student's growth\nover time with grouped bar\ncharts and detailed data tables"),
    ("Student Ranking", "Side-by-side ranking of students\nby Total Letter ID % and\nLetter Sound % with progress bars"),
    ("Cohort Progress", "Averaged metrics across selected\nstudents to view whole-group\ntrends and class-level growth"),
    ("Export PDF", "Generate printable PDF reports\nwith scores, charts, and\ncolor-coded letter grids"),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.15)
    y = Inches(2.2) + Inches(row * 2.4)

    # Card
    add_rect(slide, x, y, Inches(3.85), Inches(2.1), LIGHT_GRAY, line_color=BLUE)
    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.45),
                title, font_size=18, bold=True, color=BLUE)
    # Thin line
    add_rect(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.45), Inches(0.02), BLUE)
    # Description
    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.2),
                desc, font_size=14, color=DARK_TEXT, line_spacing=1.3)

add_speaker_notes(slide,
    "Walk through the app's key features. Emphasize that it was purpose-built for "
    "this research project. Highlight the Live Assessment (flashcard-based, three phases), "
    "the color-coded scoring system, and the ability to track individual and cohort "
    "progress over time. If you have a screenshot or live demo, show it here.")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 – Assessment Flow / How It Works
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "How It Works: Live Assessment Flow")

# Step boxes in a horizontal flow
steps = [
    ("1", "Select\nStudent", "Choose a student and\nset letter order\n(sequential or random)"),
    ("→", "", ""),
    ("2", "Uppercase\nPhase", "Flash each of 26\nuppercase letters;\nmark Identified or Not"),
    ("→", "", ""),
    ("3", "Lowercase\nPhase", "Flash each of 26\nlowercase letters;\nmark Identified or Not"),
    ("→", "", ""),
    ("4", "Sounds\nPhase", "Flash each of 26\nlowercase letters;\nstudent says the sound"),
    ("→", "", ""),
    ("5", "Review &\nSave", "See results, color-coded\nletter grids, and save\nto Google Sheets"),
]

x = Inches(0.4)
for num, title, desc in steps:
    if num == "→":
        # Arrow
        add_textbox(slide, x, Inches(3.0), Inches(0.5), Inches(0.5),
                    "→", font_size=28, color=BLUE, alignment=PP_ALIGN.CENTER,
                    bold=True)
        x += Inches(0.5)
        continue

    # Step card
    card_w = Inches(2.2)
    add_rect(slide, x, Inches(1.8), card_w, Inches(3.8), LIGHT_GRAY, line_color=BLUE)

    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), Inches(2.0),
                                   Inches(0.7), Inches(0.7))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_textbox(slide, x + Inches(0.1), Inches(2.85), Inches(2.0), Inches(0.7),
                title, font_size=16, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER, line_spacing=1.15)

    # Description
    add_textbox(slide, x + Inches(0.1), Inches(3.7), Inches(2.0), Inches(1.5),
                desc, font_size=12, color=MED_GRAY,
                alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    x += Inches(2.35)

# Color scoring legend at bottom
add_textbox(slide, Inches(0.8), Inches(6.1), Inches(3), Inches(0.4),
            "Color-Coded Scoring System:",
            font_size=15, bold=True, color=DARK_TEXT)

add_color_swatch(slide, Inches(4.2), Inches(6.15), GREEN, "≥ 80%  Proficient")
add_color_swatch(slide, Inches(6.4), Inches(6.15), YELLOW_C, "≥ 65%  Approaching")
add_color_swatch(slide, Inches(8.8), Inches(6.15), ORANGE, "≥ 40%  Developing")
add_color_swatch(slide, Inches(11.0), Inches(6.15), RED, "< 40%  Beginning")

add_speaker_notes(slide,
    "Walk through the assessment flow step by step. Emphasize: (1) sequential phases "
    "ensure complete data, (2) each letter is individually assessed, (3) results are "
    "instantly color-coded, (4) data saves directly to Google Sheets for longitudinal "
    "tracking. Point out the color-coded scoring legend at the bottom.")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 – Data Collection
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Data Collection")

# Left column
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
            "What Data Is Collected",
            font_size=24, bold=True, color=DARK_BLUE)

add_bullet_frame(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
    [
        "Per-letter identification status (Known / Unknown) for all 26 letters",
        "Three separate scores per assessment: Uppercase %, Lowercase %, Sounds %",
        "Computed Total Letter ID % (average of Uppercase + Lowercase)",
        "Timestamped assessment history for longitudinal tracking",
        "Individual letter-level data stored in Google Sheets",
        "Assessment history preserved for pre/post comparison",
    ],
    font_size=16, spacing=10)

# Right column
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.6),
            "How It Maps to the Research",
            font_size=24, bold=True, color=DARK_BLUE)

mapping = [
    ("Uppercase + Lowercase Scores", "→  Letter Recognition"),
    ("Sounds Score", "→  Letter Sounds"),
    ("Total Letter ID %", "→  Literacy Skills"),
    ("Per-Letter Known/Unknown", "→  Targeted Reteaching Data"),
    ("Pre/Post Comparison", "→  Measuring Intervention Impact"),
]

y = Inches(2.3)
for left_text, right_text in mapping:
    add_rect(slide, Inches(7), y, Inches(5.5), Inches(0.75),
             LIGHT_GRAY, line_color=BLUE)
    add_textbox(slide, Inches(7.2), y + Inches(0.05), Inches(2.5), Inches(0.65),
                left_text, font_size=14, bold=True, color=DARK_BLUE)
    add_textbox(slide, Inches(9.8), y + Inches(0.05), Inches(2.5), Inches(0.65),
                right_text, font_size=14, color=DARK_TEXT)
    y += Inches(0.85)

add_speaker_notes(slide,
    "Explain what data is collected and how it connects to your research question. "
    "The left side describes the raw data; the right side maps each data point to "
    "the research constructs. Emphasize that Total Letter ID % (average of uppercase "
    "and lowercase) maps to the operational definition of 'literacy skills.'")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 – Results (placeholder)
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Results")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
            "Pre-Assessment vs. Post-Assessment Comparison",
            font_size=24, bold=True, color=DARK_BLUE)

# Placeholder table area – left
add_rect(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(3.8),
         LIGHT_GRAY, line_color=MED_GRAY)
add_textbox(slide, Inches(0.8), Inches(3.6), Inches(5.5), Inches(1),
            "[Insert Pre/Post Data Table or\nScreenshot from Student Scorecard]",
            font_size=20, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Placeholder chart area – right
add_rect(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(3.8),
         LIGHT_GRAY, line_color=MED_GRAY)
add_textbox(slide, Inches(7), Inches(3.6), Inches(5.5), Inches(1),
            "[Insert Progress Chart or\nScreenshot from Individual Progress Tab]",
            font_size=20, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Guidance text
add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
            "Tip: Use screenshots from the Letter ID Tracker app's Student Scorecard "
            "and Individual Progress tabs",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide,
    "PLACEHOLDER SLIDE — Replace the gray boxes with actual data:\n"
    "• Left: Pre/post comparison table (screenshot from Student Scorecard or export)\n"
    "• Right: Progress chart (screenshot from Individual Progress tab)\n\n"
    "Talk through the data: What were scores at baseline? What are they now? "
    "Which students showed the most growth? Were there any surprises?")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 – Analysis & Findings (placeholder)
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Analysis & Findings")

# Left side – key findings bullets
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
            "Key Findings",
            font_size=24, bold=True, color=DARK_BLUE)

add_bullet_frame(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4),
    [
        "[Finding 1: e.g., Average cohort uppercase score increased from X% to Y%]",
        "[Finding 2: e.g., Sounds recognition showed the most/least growth]",
        "[Finding 3: e.g., X out of Y students moved from red/orange to yellow/green]",
        "[Finding 4: e.g., Letters most commonly missed were ___]",
    ],
    font_size=16, color=MED_GRAY, spacing=14)

# Right side – chart placeholder
add_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5),
         LIGHT_GRAY, line_color=MED_GRAY)
add_textbox(slide, Inches(7), Inches(3.2), Inches(5.5), Inches(1),
            "[Insert Cohort Progress Chart or\nStudent Ranking Screenshot]",
            font_size=20, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
            "Tip: Use screenshots from the Cohort Progress and Student Ranking tabs",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide,
    "PLACEHOLDER SLIDE — Replace bracketed text with your actual findings.\n"
    "• Replace the gray box with a Cohort Progress or Student Ranking screenshot.\n"
    "• Discuss what the data shows in relation to your research question.\n"
    "• Connect findings back to operational definitions (letter recognition, "
    "letter sounds, literacy skills).\n"
    "• Note any patterns: Did uppercase or lowercase improve more? "
    "Was sounds the most challenging?")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 – Reflection & Next Steps (placeholder)
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Reflection & Next Steps")

# Two columns
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
            "What I Learned",
            font_size=24, bold=True, color=DARK_BLUE)

add_bullet_frame(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.5),
    [
        "[How the UFLI sequence impacted student learning]",
        "[How the Letter ID Tracker supported data-driven instruction]",
        "[What worked well in the intervention process]",
        "[What was challenging or surprising]",
    ],
    font_size=16, color=MED_GRAY, spacing=14)

add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.6),
            "Next Steps",
            font_size=24, bold=True, color=DARK_BLUE)

add_bullet_frame(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(3.5),
    [
        "[How you would modify the intervention]",
        "[How the technology tool could be improved]",
        "[Implications for your teaching practice]",
        "[How this research informs future instruction]",
    ],
    font_size=16, color=MED_GRAY, spacing=14)

# Bottom bar
add_rect(slide, Inches(0), SLIDE_H - Inches(0.7), SLIDE_W, Inches(0.7), LIGHT_GRAY)
add_textbox(slide, Inches(0.8), SLIDE_H - Inches(0.6), Inches(11.5), Inches(0.5),
            "\"The goal is not to prove, but to improve.\" — Action Research",
            font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER,
            font_name="Calibri")

add_speaker_notes(slide,
    "PLACEHOLDER SLIDE — Replace bracketed text with your personal reflections.\n"
    "This is where you address the rubric's 'Quality of conclusion' criterion.\n"
    "• Reflect on the research question: Was it answered? Partially?\n"
    "• Discuss the role of the Letter ID Tracker as a data collection tool.\n"
    "• Be honest about limitations and what you'd change.\n"
    "• End with professional growth and how this shapes your practice going forward.")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 – Questions
# ══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Blue accent bar at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.25), BLUE)

add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
            "Questions?",
            font_size=52, bold=True, color=DARK_BLUE,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(1),
            "Thank you for your time and feedback.",
            font_size=24, color=MED_GRAY,
            alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_rect(slide, Inches(0), SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), BLUE)

add_speaker_notes(slide,
    "Open the floor for questions. The rubric includes 'Response to Questions' (2 pts) — "
    "be prepared to discuss:\n"
    "• Why you chose the UFLI sequence specifically\n"
    "• How you selected the target students\n"
    "• Technical details about the Letter ID Tracker if asked\n"
    "• Limitations of the study\n"
    "• What you would do differently")


# ── Save ─────────────────────────────────────────────────────────────

output_path = "/Users/kgath/letter-id-tracker/capstone_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
